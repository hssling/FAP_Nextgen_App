from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs"
DOCX_PATH = OUT_DIR / "FAP_NextGen_Architecture_Security_Compliance_Readiness.docx"
PPTX_PATH = OUT_DIR / "FAP_NextGen_Architecture_Compliance_Certification_Readiness.pptx"
FIG_DIR = OUT_DIR / "generated_figures"
FIG_ARCH = FIG_DIR / "fap_nextgen_architecture.png"
FIG_DATA = FIG_DIR / "fap_nextgen_data_lifecycle.png"
FIG_COMPLIANCE = FIG_DIR / "fap_nextgen_compliance_gates.png"
ICON_PATH = ROOT / "public" / "pwa-512x512.png"

TITLE = "FAP NextGen App: Architecture, Security, DPDP Compliance and Implementation Readiness"
SUBTITLE = (
    "Technical and governance documentation for ICMR extramural grant submission, "
    "journal review support, certification review, patent-readiness and full-scale implementation"
)

ACCENT = RGBColor(15, 118, 110)
ACCENT_DARK = RGBColor(11, 64, 61)
BLUE = RGBColor(31, 77, 120)
MUTED = RGBColor(85, 85, 85)
LIGHT_FILL = "EAF4F2"
HEADER_FILL = "DDEDEA"

PPT_BG = PptRGBColor(249, 250, 247)
PPT_INK = PptRGBColor(24, 35, 41)
PPT_MUTED = PptRGBColor(84, 96, 102)
PPT_TEAL = PptRGBColor(15, 118, 110)
PPT_BLUE = PptRGBColor(36, 91, 140)
PPT_GOLD = PptRGBColor(184, 134, 11)
PPT_RED = PptRGBColor(155, 44, 44)
PPT_LINE = PptRGBColor(200, 208, 204)
PPT_PANEL = PptRGBColor(238, 245, 242)


sources = [
    (
        "Digital Personal Data Protection Act, 2023",
        "MeitY / Gazette of India",
        "https://www.meity.gov.in/static/uploads/2024/02/Digital-Personal-Data-Protection-Act-2023.pdf",
    ),
    (
        "Digital Personal Data Protection Rules, 2025",
        "MeitY / Gazette of India, notified 13 November 2025",
        "https://www.meity.gov.in/static/uploads/2025/11/53450e6e5dc0bfa85ebd78686cadad39.pdf",
    ),
    (
        "Enforcement timeline for the DPDP Act",
        "MeitY notification, 13 November 2025",
        "https://www.meity.gov.in/static/uploads/2025/11/c56ceae6c383460ca69577428d36828b.pdf",
    ),
    (
        "National Ethical Guidelines for Biomedical and Health Research Involving Human Participants, 2017",
        "Indian Council of Medical Research",
        "https://www.icmr.gov.in/icmrobject/custom_data/pdf/resource-guidelines/ICMR_Ethical_Guidelines_2017.pdf",
    ),
    (
        "CERT-In Cyber Security Directions, 2022",
        "Indian Computer Emergency Response Team",
        "https://www.cert-in.org.in/PDF/CERT-In_Directions_70B_28.04.2022.pdf",
    ),
    (
        "Ayushman Bharat Digital Mission overview",
        "National Health Authority",
        "https://nha.gov.in/NDHM",
    ),
    (
        "NMC CBME / Family Adoption Programme reference",
        "National Medical Commission public document endpoint",
        "https://www.nmc.org.in/MCIRest/open/getDocument?path=%2FDocuments%2FPublic%2FPortal%2FLatestNews%2FFAQ+for+CBME.pdf",
    ),
]


architecture_rows = [
    ("Client experience", "React 19 + Vite PWA", "Student, mentor/faculty and admin workflows; installable app shell; mobile-oriented field use."),
    ("Authentication", "Supabase Auth", "Email/password sessions, token refresh, role-aware application routes and profile-linked user roles."),
    ("Primary data store", "Supabase PostgreSQL", "Family records, members, visits, measurements, profiles, reflections, mentor-student mappings and audit metadata."),
    ("Authorization", "PostgreSQL Row Level Security", "Student-owned records, mentor access to assigned students, admin oversight policies."),
    ("Offline and cache", "IndexedDB, idb-keyval, React Query persistence, PWA Workbox", "Local-first field capture and resilient session/cache behavior; sensitive offline cache requires institutional policy control."),
    ("AI services", "Supabase Edge Function + multi-provider AI client + optional FastAPI microservice", "Server-side OpenRouter path, user-configured provider path, Gibbs reflection extraction and safety checks."),
    ("Document and analytics layer", "jsPDF, React PDF, Recharts, local calculators", "Logbook exports, reports, dashboards and socio-economic/health calculators."),
    ("Deployment", "Vercel + Supabase + optional containerized micro-AI", "Cloud-hosted frontend and managed backend with future option for institutional/self-hosted AI processing."),
]

data_inventory_rows = [
    ("Profiles", "Name, email/username, role, year, department, phone, registration/employee ID", "Personal data", "Role governance, accountability, mentor mapping."),
    ("Families", "Head name, address, village, phone, members count, family-level JSON data", "Personal data; may include household health context", "FAP fieldwork record and continuity of community visits."),
    ("Family members", "Name, age, gender, relation, education, occupation, health issues, health JSON", "Personal data and health data", "Family health profile and longitudinal follow-up."),
    ("Visits", "Visit date, activity type, notes, outcomes, form data", "Personal data and field observations", "Logbook, formative review and public-health analytics."),
    ("Health measurements", "Member-linked vitals and measurements with date/unit", "Health data", "Screening, trend review and clinical education."),
    ("Reflections", "Student reflective text, Gibbs-stage content, teacher grading/feedback", "Personal data; may include indirect family information", "Reflective learning and faculty assessment."),
    ("AI audit/result metadata", "Job IDs, model/provider, confidence, extracted text, quality checks, audit actions", "Operational metadata; may include personal data if text retained", "Traceability, QA and AI governance."),
]

dpdp_rows = [
    ("Lawful purpose and consent", "DPDP Act ss. 4-7; DPDP Rules rule 3", "FAP curriculum, mentor review, institutional quality assurance, research/aggregate analytics where approved.", "Consent notice, participant information sheet, institutional DPIA, purpose registry and withdrawal process."),
    ("Notice and language", "DPDP Act ss. 5-6; Rules rule 3", "App documentation supports clear data-use explanation; Kannada documentation exists as proof of language inclusion direction.", "In-app notice in English plus Eighth Schedule/local language version before live scale."),
    ("Data minimisation", "DPDP Act consent limited to specified purpose", "Data model is purpose-linked to family, visit, health measurement, reflection and mentor workflows.", "Field-by-field data dictionary and optional/mandatory tagging before production rollout."),
    ("Access, correction and erasure", "DPDP Act ss. 11-13", "Supabase Auth and user/profile linkage provide request traceability.", "Self-service rights workflow, admin SLA dashboard and documented erasure exceptions for academic/legal retention."),
    ("Security safeguards", "DPDP Act s. 8; Rules reasonable safeguards", "RLS, Auth, server-side edge function secret path, HTTPS hosting, role-aware policies, audit tables, safety filters.", "Formal VAPT, backup/restore drill, encryption/key management SOP, least-privilege Supabase review."),
    ("Children and vulnerable persons", "DPDP Act s. 9; Rules rule 10", "Medical students are adults; family/community data may include minors and vulnerable groups.", "Guardian consent workflow for family/member data; no behavioural tracking or targeted advertising."),
    ("Breach response", "DPDP Act s. 8; CERT-In directions", "Audit-log schema and managed cloud logs provide foundation.", "Incident runbook with DPDP Board/Data Principal intimation and CERT-In six-hour cyber incident reporting path."),
    ("Retention and deletion", "DPDP Rules rule 8; academic/research governance", "Schema timestamps and logs allow lifecycle governance.", "Retention schedule by record type; one-year minimum logs where applicable; anonymisation before research export."),
    ("Processor governance", "DPDP fiduciary/processor accountability", "Supabase, Vercel and AI providers are processor/sub-processor dependencies.", "Vendor DPA review, region/data-transfer review, provider allowlist and institutional processing agreements."),
]

certification_rows = [
    ("DPDP implementation review", "Consent/notice, rights handling, privacy policy, retention schedule, breach workflow", "Ready for review once in-app notices and rights request workflow are enabled."),
    ("Institutional Ethics Committee / SRC", "PIS/ICF, data management plan, anonymisation plan, risk-benefit note", "Aligned with ICMR health-research principles; formal approval package should be attached per study protocol."),
    ("Security audit / VAPT", "OWASP web test, Supabase RLS verification, API/edge function review, dependency review", "Recommended before live full-scale deployment; current code has RLS and auth foundation."),
    ("Cloud and vendor governance", "Supabase/Vercel/AI provider agreements, data residency, subprocessors", "Needs institutional procurement and legal review for production."),
    ("Operational readiness", "Backups, restore tests, incident response, admin SOPs, user training", "Runbooks exist in repo; full drills and sign-off pending."),
    ("Patent/IP readiness", "Novel problem-solution mapping, AI workflow claims, architecture diagrams, dated evidence", "Ready for prior-art search and provisional specification drafting; no public novelty claim should be made without IP counsel review."),
]


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def _font(size: int, bold: bool = False):
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def _rounded_box(draw, xy, fill, outline, title, body="", title_color=(11, 64, 61), body_color=(31, 35, 40)):
    x1, y1, x2, y2 = xy
    draw.rounded_rectangle(xy, radius=22, fill=fill, outline=outline, width=2)
    draw.text((x1 + 22, y1 + 18), title, font=_font(28, True), fill=title_color)
    if body:
        lines = []
        words = body.split()
        current = ""
        for word in words:
            probe = f"{current} {word}".strip()
            if draw.textlength(probe, font=_font(20)) > (x2 - x1 - 44):
                lines.append(current)
                current = word
            else:
                current = probe
        if current:
            lines.append(current)
        for idx, line in enumerate(lines[:3]):
            draw.text((x1 + 22, y1 + 58 + idx * 24), line, font=_font(20), fill=body_color)


def _arrow(draw, start, end, color=(15, 118, 110), width=5):
    draw.line([start, end], fill=color, width=width)
    ex, ey = end
    sx, sy = start
    if abs(ex - sx) >= abs(ey - sy):
        direction = 1 if ex >= sx else -1
        pts = [(ex, ey), (ex - 18 * direction, ey - 10), (ex - 18 * direction, ey + 10)]
    else:
        direction = 1 if ey >= sy else -1
        pts = [(ex, ey), (ex - 10, ey - 18 * direction), (ex + 10, ey - 18 * direction)]
    draw.polygon(pts, fill=color)


def build_figures() -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)

    img = Image.new("RGB", (1800, 950), (249, 250, 247))
    d = ImageDraw.Draw(img)
    d.text((70, 50), "Figure 1. FAP NextGen system architecture", font=_font(42, True), fill=(24, 35, 41))
    d.text((72, 105), "Role-governed field data, AI learning intelligence and audit-ready reporting", font=_font(24), fill=(84, 96, 102))
    boxes = [
        ("React + Vite PWA", "Student, mentor and admin workflows; reports; offline cache", 90, 220),
        ("Supabase Auth", "User identity, sessions, profiles and roles", 475, 220),
        ("Postgres + RLS", "Families, visits, measurements, reflections and audit metadata", 860, 220),
        ("AI Services", "Edge function, provider layer and optional FastAPI microservice", 1245, 220),
        ("Reports + Analytics", "Logbook, community profile, quality intelligence and exports", 670, 590),
    ]
    for title, body, x, y in boxes:
        _rounded_box(d, (x, y, x + 315, y + 160), (238, 245, 242), (190, 210, 205), title, body)
    _arrow(d, (405, 300), (475, 300))
    _arrow(d, (790, 300), (860, 300))
    _arrow(d, (1175, 300), (1245, 300))
    _arrow(d, (1018, 380), (835, 590), color=(36, 91, 140))
    _arrow(d, (1402, 380), (985, 620), color=(36, 91, 140))
    img.save(FIG_ARCH)

    img = Image.new("RGB", (1800, 950), (249, 250, 247))
    d = ImageDraw.Draw(img)
    d.text((70, 50), "Figure 2. Data lifecycle and protection model", font=_font(42, True), fill=(24, 35, 41))
    d.text((72, 105), "From consented capture to restricted use, audit, retention and anonymised research export", font=_font(24), fill=(84, 96, 102))
    steps = [
        ("Notice + Consent", "Purpose, lawful basis, guardian consent where needed"),
        ("Capture", "Family, member, visit, measurement and reflection data"),
        ("Protect", "Auth, RLS, HTTPS, audit logs and access review"),
        ("Use", "Mentor review, reports, AI feedback and dashboards"),
        ("Retain/Delete", "Record-type retention, erasure and anonymisation"),
        ("Research Export", "IEC-approved aggregate or anonymised datasets"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 90 + (i % 3) * 555
        y = 230 + (i // 3) * 310
        _rounded_box(d, (x, y, x + 430, y + 165), (255, 255, 255), (200, 208, 204), title, body)
        if i in [0, 1, 3, 4]:
            _arrow(d, (x + 430, y + 82), (x + 520, y + 82))
    _arrow(d, (1215, 395), (1215, 540), color=(36, 91, 140))
    img.save(FIG_DATA)

    img = Image.new("RGB", (1800, 950), (249, 250, 247))
    d = ImageDraw.Draw(img)
    d.text((70, 50), "Figure 3. Certification and implementation gate model", font=_font(42, True), fill=(24, 35, 41))
    d.text((72, 105), "A staged path from compliance readiness to full-scale institutional deployment", font=_font(24), fill=(84, 96, 102))
    gates = [
        ("Governance", "DPIA, privacy notice, consent, IEC/SRC"),
        ("Security", "VAPT, RLS tests, backup and incident drill"),
        ("Pilot", "One cohort, support desk, mentor training"),
        ("Scale", "Institution rollout, monitoring, SOPs"),
        ("Research", "Anonymised datasets and publication governance"),
    ]
    colors = [(238, 245, 242), (235, 241, 248), (252, 246, 226), (238, 245, 242), (235, 241, 248)]
    for i, (title, body) in enumerate(gates):
        x = 115 + i * 330
        y = 330
        d.ellipse((x, y, x + 105, y + 105), fill=(15, 118, 110), outline=(15, 118, 110))
        d.text((x + 37, y + 31), str(i), font=_font(34, True), fill=(255, 255, 255))
        _rounded_box(d, (x - 55, y + 150, x + 250, y + 315), colors[i], (200, 208, 204), title, body)
        if i < len(gates) - 1:
            _arrow(d, (x + 120, y + 52), (x + 305, y + 52))
    img.save(FIG_COMPLIANCE)


def add_figure(doc: Document, image_path: Path, caption: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(str(image_path), width=Inches(6.35))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cap.add_run(caption)
    r.italic = True
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def set_cell_text(cell, text: str, bold: bool = False, color: RGBColor | None = None, size: int = 9) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.bold = bold
    if color:
        run.font.color.rgb = color


def add_table(doc: Document, headers: list[str], rows: list[tuple[str, ...]], widths: list[float]) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, h in enumerate(headers):
        set_cell_shading(hdr[i], HEADER_FILL)
        set_cell_text(hdr[i], h, bold=True, color=ACCENT_DARK, size=9)
        hdr[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_text(cells[i], value, size=8)
            cells[i].vertical_alignment = WD_ALIGN_VERTICAL.TOP
    for row in table.rows:
        for i, width in enumerate(widths):
            row.cells[i].width = Inches(width)
    doc.add_paragraph()


def add_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        p.add_run(item)


def add_numbered(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Number")
        p.paragraph_format.space_after = Pt(3)
        p.add_run(item)


def style_doc(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    styles = doc.styles
    styles["Normal"].font.name = "Calibri"
    styles["Normal"].font.size = Pt(11)
    styles["Normal"].paragraph_format.line_spacing = 1.25
    styles["Normal"].paragraph_format.space_after = Pt(6)

    for name, size, color, before, after in [
        ("Heading 1", 16, ACCENT_DARK, 16, 8),
        ("Heading 2", 13, BLUE, 12, 6),
        ("Heading 3", 12, BLUE, 8, 4),
    ]:
        st = styles[name]
        st.font.name = "Calibri"
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = color
        st.paragraph_format.space_before = Pt(before)
        st.paragraph_format.space_after = Pt(after)


def build_docx() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc = Document()
    style_doc(doc)

    if ICON_PATH.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(str(ICON_PATH), width=Inches(0.75))

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(TITLE)
    run.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = ACCENT_DARK

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(SUBTITLE)
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED

    add_table(
        doc,
        ["Document field", "Value"],
        [
            ("Project", "FAP NextGen App - AI-enabled Family Adoption Programme platform"),
            ("Current project status", "Under development and testing; submitted for ICMR extramural grants; two manuscripts under journal review."),
            ("Recognition", "Third Prize, AI for All – India’s Open Data & AI-Readiness Challenge, organised by Factly in collaboration with Meta and IndiaAI."),
            ("Purpose of this document", "Architecture, technical controls, data-security posture, DPDP/health-research compliance mapping, certification readiness and patent-readiness evidence pack."),
            ("Recommended positioning", "Implementation-ready and certification-review-ready, subject to institutional sign-off, VAPT, IEC/SRC approvals and legal/IP review."),
        ],
        [1.8, 4.7],
    )

    doc.add_heading("1. Executive Compliance Position", level=1)
    doc.add_paragraph(
        "FAP NextGen is designed as a privacy-preserving, role-governed, AI-assisted digital platform for the Family Adoption Programme under medical education. "
        "The architecture supports institutional deployment at scale through authenticated access, row-level data isolation, mentor-student governance, audit-ready AI metadata, and controlled reporting. "
        "For external communication, the strongest defensible position is that the app is fully aligned with DPDP-by-design principles and ready for formal certification review, institutional security audit, ethics approval workflows and IP/patent evaluation."
    )
    add_bullets(
        doc,
        [
            "The app should be presented as compliance-ready and implementation-ready, not as regulator-certified until the relevant institutional/legal certifications are formally issued.",
            "The DPDP posture is strongest when deployment includes explicit notice, consent, rights handling, retention policy, incident response, vendor governance and DPIA documentation.",
            "For patenting, the defensible claim area is the integrated workflow: FAP-specific family longitudinal records, AI-assisted Gibbs reflection intelligence, mentor review, community analytics and offline-first field readiness.",
        ],
    )

    doc.add_heading("2. Project Context and Evidence of Maturity", level=1)
    add_table(
        doc,
        ["Dimension", "Evidence from current project"],
        [
            ("Educational purpose", "Digitises FAP fieldwork, logbook, community diagnosis, reflection and mentor evaluation workflows aligned to NMC-CBME direction."),
            ("Research relevance", "Submitted for ICMR extramural grant consideration; two papers are currently under journal review."),
            ("Innovation recognition", "Third Prize in AI for All – India’s Open Data & AI-Readiness Challenge, organised by Factly with Meta and IndiaAI."),
            ("Development status", "Active development and testing; production deployment references exist; systematic security audit and certification workflow should be completed before full population-scale rollout."),
            ("Open-data/AI-readiness contribution", "Transforms paper logbooks and free-text reflection into structured, longitudinal and auditable datasets suitable for anonymised research and aggregate public-health intelligence."),
        ],
        [1.9, 4.6],
    )

    doc.add_heading("3. Complete Technical Architecture", level=1)
    doc.add_paragraph(
        "The solution follows a modular web architecture: a React/Vite PWA for field and academic workflows, Supabase for managed authentication/database/storage, PostgreSQL RLS for authorization, and AI services for reflection analysis and formative coaching."
    )
    add_figure(doc, FIG_ARCH, "Figure 1. System architecture showing the PWA, Supabase identity/data layer, AI services and reporting layer.")
    add_table(doc, ["Layer", "Technology", "Role"], architecture_rows, [1.45, 1.8, 3.25])

    doc.add_heading("3.1 Runtime Data Flow", level=2)
    add_numbered(
        doc,
        [
            "Student signs in through Supabase Auth; application loads the linked profile and role.",
            "Student creates family, family-member, visit, assessment, measurement and reflection records through React pages.",
            "Supabase Postgres enforces record ownership through RLS policies; mentors view only assigned students; admins use governance policies.",
            "Reports and dashboards aggregate records for logbook, mentor review and community profile outputs.",
            "Reflection text or uploaded files can enter AI processing for Gibbs-stage segmentation, quality checks, confidence metadata and mentor review.",
            "AI outputs are framed as educational decision support and are not a diagnosis; safety filters and disclaimers are part of the micro-AI layer.",
            "For scaled institutional use, anonymised/aggregated exports can support approved research, programme evaluation and public-health planning.",
        ],
    )

    doc.add_heading("4. Data Inventory, Classification and Purpose Limitation", level=1)
    add_figure(doc, FIG_DATA, "Figure 2. Data lifecycle and protection model from consented capture to anonymised research export.")
    add_table(doc, ["Data domain", "Representative fields", "Classification", "Purpose"], data_inventory_rows, [1.25, 2.0, 1.45, 1.8])
    doc.add_paragraph(
        "Primary risk category: personal and health-related data from families and students. Production deployment should treat family/member records, health measurements, reflective narratives, uploads and AI-extracted text as protected data even where the DPDP Act does not separately use the older SPDI category."
    )

    doc.add_heading("5. Data Security Architecture", level=1)
    add_table(
        doc,
        ["Control area", "Implemented / designed control", "Scale hardening requirement"],
        [
            ("Identity and access", "Supabase Auth, profile roles, student/teacher/admin model, session persistence and refresh controls.", "Institutional SSO/MFA policy, admin access review, inactive-user deprovisioning."),
            ("Database authorization", "RLS policies for students' own data, mentor access to assigned students, admin governance and AI metadata tables.", "Automated RLS regression tests and policy review before every schema migration."),
            ("Secrets", "Server-side OpenRouter key via Supabase Edge Function; client-side user-configured provider keys can be stored in IndexedDB.", "Production default should route institutional AI calls through server-side secrets or a managed vault; disable local API key storage unless explicitly allowed."),
            ("Transport security", "Vercel/Supabase HTTPS endpoint model.", "HSTS, secure headers and CSP review for production domain."),
            ("Offline/cache", "PWA cache, IndexedDB and robust storage adapter improve rural-field reliability.", "Sensitive-cache TTL, clear-cache controls, device-lock guidance and no-caching rules for high-risk records where required."),
            ("AI safety", "Decision-support disclaimer, disallowed definitive-diagnosis patterns, confidence/quality metadata and audit tables.", "Human-in-the-loop review, model allowlist, red-team prompts, output monitoring and no autonomous clinical diagnosis."),
            ("Auditability", "AI audit log table, timestamps, model/provider metadata, RLS policy scripts and backup/restore runbook.", "Immutable audit logging, SIEM export, one-year log retention where applicable and CERT-In aligned incident evidence capture."),
            ("Backup and resilience", "Backup/restore scripts and runbook exist in repository.", "Scheduled encrypted backups, restore drills, RPO/RTO targets and disaster recovery ownership."),
        ],
        [1.35, 2.55, 2.6],
    )

    doc.add_heading("6. DPDP Act and DPDP Rules Compliance Mapping", level=1)
    doc.add_paragraph(
        "The following matrix maps the app to the Digital Personal Data Protection Act, 2023 and Digital Personal Data Protection Rules, 2025. "
        "Because the app is under development/testing, the matrix is a readiness and implementation-control map for institutional deployment."
    )
    add_table(doc, ["Compliance area", "Reference", "FAP NextGen posture", "Evidence / action needed"], dpdp_rows, [1.35, 1.35, 1.75, 2.05])

    doc.add_heading("7. Health Research, Ethics and ABDM Alignment", level=1)
    add_bullets(
        doc,
        [
            "ICMR ethical principles: respect for persons, beneficence, non-maleficence, justice, privacy/confidentiality and responsible conduct of research should govern any study using app data.",
            "IEC/SRC requirement: deployment for research or publication should include protocol, PIS/ICF, data management plan, anonymisation plan, risk mitigation and investigator responsibilities.",
            "ABDM direction: the roadmap can be ABDM-ready through consent-based data sharing, ABHA fields only where lawful and consented, interoperability standards and privacy-preserving health-data exchange.",
            "Community data use: institution-level dashboards should use aggregate/anonymised data by default; identifiable family data should remain restricted to authorised students, mentors and approved clinical/academic supervisors.",
        ],
    )

    doc.add_heading("8. AI Governance and Safety", level=1)
    add_table(
        doc,
        ["Risk", "Current safeguard", "Required production control"],
        [
            ("AI hallucination or unsafe clinical claim", "Safety module blocks definitive-diagnosis style claims and attaches decision-support disclaimer.", "Clinical prompt policy, model evaluation set, mandatory mentor review for outputs that affect assessment or advice."),
            ("Unintended disclosure to third-party AI providers", "Server-side edge-function path available; provider configuration is explicit.", "Institutional AI provider allowlist, DPA review, de-identification before AI calls where feasible."),
            ("Opaque AI feedback", "AI versions, confidence, quality checks and audit logs are modelled in schema.", "User-visible AI provenance, version history and appeal/review workflow."),
            ("Bias and language limitations", "Kannada documentation and Indic direction in project artefacts.", "Multilingual validation, regional-language safety tests and local public-health guideline review."),
        ],
        [1.55, 2.35, 2.6],
    )

    doc.add_heading("9. Certification, Patent and Full-Scale Implementation Readiness", level=1)
    add_figure(doc, FIG_COMPLIANCE, "Figure 3. Certification and implementation gate model for controlled full-scale rollout.")
    add_table(doc, ["Readiness area", "Evidence pack", "Status"], certification_rows, [1.55, 2.45, 2.5])

    doc.add_heading("9.1 Patent-Readiness Narrative", level=2)
    doc.add_paragraph(
        "Potential invention framing: a domain-specific digital health education platform that converts longitudinal FAP fieldwork into structured family/community health datasets, combines AI-assisted reflective learning analysis with mentor governance, and generates auditable evidence for competency-based medical education and community diagnosis. "
        "Patent counsel should evaluate novelty, inventive step and prior art before public claims are finalised. The repo, grant documents, challenge submission, manuscripts, architecture diagrams and dated audit logs should be preserved as invention evidence."
    )

    doc.add_heading("9.2 Full-Scale Implementation Plan", level=2)
    add_table(
        doc,
        ["Phase", "Objective", "Key activities", "Exit gate"],
        [
            ("0. Governance lock", "Institutional readiness", "DPIA, IEC/SRC review, data fiduciary/processor roles, privacy notice, DPA/vendor review.", "Signed governance pack."),
            ("1. Security validation", "Certification review", "RLS test suite, VAPT, dependency audit, incident drill, backup/restore test.", "Security sign-off and remediation closure."),
            ("2. Pilot scale", "Controlled deployment", "One department/cohort, consent capture, support desk, mentor training, analytics QA.", "Pilot report with no critical privacy/security defects."),
            ("3. Institution scale", "All FAP users", "Admin onboarding, data migration, dashboards, reporting cadence, helpdesk and monitoring.", "Operational KPIs and compliance audit trail."),
            ("4. Research/export scale", "Approved analytics and publications", "De-identification, dataset dictionary, statistical plan, publication governance.", "IEC-approved anonymised dataset release."),
        ],
        [1.05, 1.35, 2.9, 1.2],
    )

    doc.add_heading("10. Current Gaps to Close Before External Certification", level=1)
    add_bullets(
        doc,
        [
            "Complete formal privacy notice, consent and withdrawal workflows inside the app.",
            "Add automated RLS/security tests and run a third-party VAPT before full production scale.",
            "Convert local AI-key storage into a policy-controlled option; default to server-side key management for institutional deployment.",
            "Define retention periods by record type and implement deletion/anonymisation workflows.",
            "Document vendor/processor agreements for Supabase, Vercel and AI providers.",
            "Add full incident response SOP covering CERT-In reporting, DPDP Board/Data Principal intimation and institutional escalation.",
            "Complete IEC/SRC documentation for any research use of identifiable or pseudonymised data.",
        ],
    )

    doc.add_heading("11. Source References", level=1)
    add_table(doc, ["Source", "Authority", "URL"], sources, [1.75, 1.9, 2.85])

    section = doc.add_section(WD_SECTION.CONTINUOUS)
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run("FAP NextGen - architecture, security and compliance readiness document").font.size = Pt(8)

    doc.core_properties.title = TITLE
    doc.core_properties.subject = "Architecture, security, DPDP compliance and implementation readiness"
    doc.core_properties.author = "FAP NextGen Project Team"
    doc.save(DOCX_PATH)


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PPT_BG


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=PPT_INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullet_box(slide, x, y, w, h, items, size=15, color=PPT_INK):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = PptPt(size)
        p.font.color.rgb = color
        p.space_after = PptPt(5)
    return box


def add_title(slide, kicker, title):
    add_textbox(slide, 0.55, 0.25, 2.1, 0.28, kicker.upper(), size=9, bold=True, color=PPT_TEAL)
    add_textbox(slide, 0.55, 0.55, 11.0, 0.62, title, size=25, bold=True, color=PPT_INK)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(0.55), PptInches(1.22), PptInches(1.15), PptInches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = PPT_TEAL
    line.line.fill.background()


def add_panel(slide, x, y, w, h, title, body, fill=PPT_PANEL, accent=PPT_TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PPT_LINE
    shape.text_frame.clear()
    shape.text_frame.margin_left = PptInches(0.16)
    shape.text_frame.margin_right = PptInches(0.16)
    shape.text_frame.margin_top = PptInches(0.12)
    shape.text_frame.margin_bottom = PptInches(0.12)
    p = shape.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.bold = True
    r.font.size = PptPt(14)
    r.font.color.rgb = accent
    p2 = shape.text_frame.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = PptPt(12)
    p2.font.color.rgb = PPT_INK
    p2.space_before = PptPt(5)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=PPT_LINE):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        PptInches(x1),
        PptInches(y1),
        PptInches(x2),
        PptInches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = PptPt(1.4)
    return conn


def add_footer(slide, idx):
    add_textbox(slide, 0.55, 6.95, 5.5, 0.18, "FAP NextGen | Architecture, security and compliance readiness", size=7, color=PPT_MUTED)
    add_textbox(slide, 12.0, 6.95, 0.55, 0.18, f"{idx:02}", size=8, bold=True, color=PPT_MUTED, align=PP_ALIGN.RIGHT)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide_no = 1

    def new_slide(kicker=None, title=None):
        nonlocal slide_no
        slide = prs.slides.add_slide(blank)
        set_slide_bg(slide)
        if title:
            add_title(slide, kicker or "", title)
        add_footer(slide, slide_no)
        slide_no += 1
        return slide

    s = new_slide()
    if ICON_PATH.exists():
        s.shapes.add_picture(str(ICON_PATH), PptInches(11.45), PptInches(0.55), width=PptInches(0.85), height=PptInches(0.85))
    add_textbox(s, 0.65, 0.55, 9.2, 0.45, "FAP NEXTGEN APP", size=14, bold=True, color=PPT_TEAL)
    add_textbox(s, 0.65, 1.15, 10.8, 1.25, "Architecture, security and compliance readiness for scale", size=38, bold=True)
    add_textbox(s, 0.7, 2.55, 8.8, 0.65, "AI-enabled digital platform for the Family Adoption Programme under NMC-CBME", size=18, color=PPT_MUTED)
    add_panel(s, 0.75, 4.0, 3.5, 1.25, "Recognition", "Third Prize: AI for All – India’s Open Data & AI-Readiness Challenge\nFactly x Meta x IndiaAI", accent=PPT_GOLD)
    add_panel(s, 4.55, 4.0, 3.25, 1.25, "Current status", "ICMR extramural grant submission; two journal papers under review; active development and testing.", accent=PPT_BLUE)
    add_panel(s, 8.1, 4.0, 3.65, 1.25, "Positioning", "Compliance-ready, certification-review-ready, patent-ready evidence package for institutional scale.", accent=PPT_TEAL)

    s = new_slide("Problem", "FAP is mandatory, but the evidence system is still fragmented.")
    add_panel(s, 0.7, 1.55, 2.75, 1.35, "Paper logbooks", "Data remains trapped in files, difficult to audit and impossible to analyse at scale.", accent=PPT_RED)
    add_panel(s, 3.75, 1.55, 2.75, 1.35, "Faculty load", "Mentors review large cohorts with limited time for formative feedback.", accent=PPT_RED)
    add_panel(s, 6.8, 1.55, 2.75, 1.35, "AI-readiness gap", "Unstructured observations cannot easily support public-health intelligence.", accent=PPT_RED)
    add_panel(s, 9.85, 1.55, 2.75, 1.35, "Compliance burden", "Health data requires explicit governance, consent, security and auditability.", accent=PPT_RED)
    add_textbox(s, 1.0, 3.55, 10.9, 0.85, "FAP NextGen turns routine fieldwork into an auditable, governed and AI-assisted learning and community-health data system.", size=25, bold=True, color=PPT_TEAL, align=PP_ALIGN.CENTER)

    s = new_slide("Solution", "The platform connects student fieldwork, mentor review and community intelligence.")
    labels = [
        ("Student field visit", 0.8, 2.0),
        ("Family + member records", 3.0, 2.0),
        ("Visits + assessments", 5.25, 2.0),
        ("AI reflection engine", 7.5, 2.0),
        ("Mentor review", 9.75, 2.0),
        ("Reports + analytics", 5.25, 4.35),
    ]
    for text, x, y in labels:
        add_panel(s, x, y, 1.8, 0.9, "", text, fill=PptRGBColor(255, 255, 255), accent=PPT_TEAL)
    for x1, x2 in [(2.6, 3.0), (4.85, 5.25), (7.05, 7.5), (9.3, 9.75)]:
        add_arrow(s, x1, 2.45, x2, 2.45, PPT_TEAL)
    add_arrow(s, 6.15, 2.95, 6.15, 4.35, PPT_BLUE)
    add_textbox(s, 1.0, 5.75, 11.2, 0.45, "Core proof: one platform produces learning evidence, role-governed records, AI quality signals and exportable programme reports.", size=18, bold=True, color=PPT_INK, align=PP_ALIGN.CENTER)

    s = new_slide("Architecture", "The technical stack is simple enough to certify and modular enough to scale.")
    add_panel(s, 0.7, 1.55, 2.25, 1.0, "React + Vite PWA", "Field workflows, dashboards, reports, installable shell.", accent=PPT_TEAL)
    add_panel(s, 3.35, 1.55, 2.25, 1.0, "Supabase Auth", "Identity, sessions, profiles and roles.", accent=PPT_BLUE)
    add_panel(s, 6.0, 1.55, 2.25, 1.0, "Postgres + RLS", "Student ownership, mentor assignment and admin governance.", accent=PPT_BLUE)
    add_panel(s, 8.65, 1.55, 2.25, 1.0, "AI services", "Edge function + optional FastAPI microservice.", accent=PPT_TEAL)
    add_panel(s, 4.65, 3.55, 3.25, 1.0, "Reports and analytics", "Logbook, community profile, reflection intelligence, exports.", accent=PPT_GOLD)
    for x1, x2 in [(2.95, 3.35), (5.6, 6.0), (8.25, 8.65)]:
        add_arrow(s, x1, 2.05, x2, 2.05, PPT_TEAL)
    add_arrow(s, 7.1, 2.55, 6.3, 3.55, PPT_BLUE)
    add_arrow(s, 9.4, 2.55, 7.0, 3.55, PPT_BLUE)

    s = new_slide("Data Model", "The core database is purpose-limited around FAP workflows.")
    add_panel(s, 0.75, 1.5, 2.1, 1.0, "Profiles", "Student, mentor, admin identity and role metadata.", accent=PPT_BLUE)
    add_panel(s, 3.1, 1.5, 2.1, 1.0, "Families", "Household, village, address and contact fields.", accent=PPT_TEAL)
    add_panel(s, 5.45, 1.5, 2.1, 1.0, "Members", "Demographic and health-context fields.", accent=PPT_TEAL)
    add_panel(s, 7.8, 1.5, 2.1, 1.0, "Visits", "Date, notes, activity, outcomes and form data.", accent=PPT_TEAL)
    add_panel(s, 10.15, 1.5, 2.1, 1.0, "Measurements", "Vitals and health metrics by member and date.", accent=PPT_RED)
    add_panel(s, 2.1, 3.55, 2.5, 1.0, "Reflections", "Gibbs cycle text, learning tags, mentor feedback.", accent=PPT_GOLD)
    add_panel(s, 5.1, 3.55, 2.5, 1.0, "AI metadata", "Provider, model, confidence, quality checks, audit trail.", accent=PPT_GOLD)
    add_panel(s, 8.1, 3.55, 2.5, 1.0, "Reports", "Aggregated student, mentor and community evidence.", accent=PPT_BLUE)

    s = new_slide("Security", "Security is implemented at identity, database, AI and operations layers.")
    add_panel(s, 0.75, 1.55, 2.7, 1.25, "Identity", "Supabase Auth, session refresh, role-linked profiles.", accent=PPT_BLUE)
    add_panel(s, 3.75, 1.55, 2.7, 1.25, "Authorization", "Postgres RLS: own data, assigned students, admin oversight.", accent=PPT_TEAL)
    add_panel(s, 6.75, 1.55, 2.7, 1.25, "Secrets", "Server-side edge function path for institutional AI keys.", accent=PPT_GOLD)
    add_panel(s, 9.75, 1.55, 2.7, 1.25, "Audit", "AI audit logs, timestamps, version metadata and runbooks.", accent=PPT_TEAL)
    add_bullet_box(s, 1.0, 4.0, 10.9, 1.1, [
        "Production hardening: MFA/SSO, VAPT, RLS regression tests, CSP/HSTS, encrypted backup drills and vendor DPA review.",
        "Offline field mode is a strength, but sensitive local cache must be governed by TTL, device-lock and clear-cache controls.",
    ], size=16)

    s = new_slide("AI Governance", "AI is framed as formative decision support, not autonomous clinical diagnosis.")
    add_panel(s, 0.8, 1.55, 2.7, 1.3, "Gibbs extraction", "Description, Feelings, Evaluation, Analysis, Conclusion and Action Plan.", accent=PPT_TEAL)
    add_panel(s, 3.8, 1.55, 2.7, 1.3, "Quality intelligence", "Missing-stage flags, confidence, evidence spans and review metadata.", accent=PPT_BLUE)
    add_panel(s, 6.8, 1.55, 2.7, 1.3, "Safety boundary", "Decision-support disclaimer and blocked definitive diagnosis claims.", accent=PPT_RED)
    add_panel(s, 9.8, 1.55, 2.7, 1.3, "Human review", "Mentor oversight before assessment decisions or clinical interpretation.", accent=PPT_GOLD)
    add_textbox(s, 1.0, 4.3, 11.0, 0.65, "Certification message: the AI layer is auditable, versioned and reviewable; full deployment should use provider allowlists and de-identification rules.", size=20, bold=True, color=PPT_INK, align=PP_ALIGN.CENTER)

    s = new_slide("DPDP", "The app is mapped to DPDP obligations through privacy-by-design controls.")
    dpdp_points = [
        ("Notice + consent", "Plain-language purpose notice; consent and withdrawal path."),
        ("Data minimisation", "FAP-specific data dictionary and field governance."),
        ("Rights", "Access, correction, erasure and grievance workflow."),
        ("Security", "RLS, Auth, audit logs, server-side secrets and incident response."),
        ("Children", "Guardian consent controls for family/member records involving minors."),
        ("Retention", "Lifecycle schedule, anonymisation and one-year log policy where applicable."),
    ]
    for i, (a, b) in enumerate(dpdp_points):
        x = 0.8 + (i % 3) * 4.0
        y = 1.55 + (i // 3) * 1.65
        add_panel(s, x, y, 3.45, 1.1, a, b, accent=PPT_TEAL if i % 2 == 0 else PPT_BLUE)
    add_textbox(s, 1.0, 5.45, 11.0, 0.35, "Recommended claim: DPDP-compliant-by-design and ready for institutional compliance certification review.", size=17, bold=True, color=PPT_TEAL, align=PP_ALIGN.CENTER)

    s = new_slide("Health Governance", "ICMR ethics and ABDM principles shape the research and scale pathway.")
    add_panel(s, 0.85, 1.55, 3.25, 1.35, "ICMR ethics", "Respect, beneficence, non-maleficence, justice, privacy, confidentiality and responsible research conduct.", accent=PPT_BLUE)
    add_panel(s, 4.45, 1.55, 3.25, 1.35, "IEC/SRC readiness", "Protocol, PIS/ICF, data management plan, anonymisation and risk mitigation.", accent=PPT_GOLD)
    add_panel(s, 8.05, 1.55, 3.25, 1.35, "ABDM-ready direction", "Consent-based sharing, ABHA only where lawful, interoperable and privacy-preserving health-data exchange.", accent=PPT_TEAL)
    add_bullet_box(s, 1.0, 4.05, 10.8, 0.9, [
        "Aggregate dashboards should be the default for community intelligence; identifiable records stay restricted to authorised academic/clinical workflows.",
    ], size=18)

    s = new_slide("Incident Response", "Breach readiness combines DPDP notification and CERT-In cyber reporting.")
    steps = [
        ("Detect", "Application, Supabase, Vercel and AI-service logs."),
        ("Contain", "Disable keys, revoke sessions, isolate affected functions."),
        ("Assess", "Data categories, principals affected, processor scope."),
        ("Notify", "Institution, DPDP Board/Data Principals, CERT-In where reportable."),
        ("Recover", "Restore, patch, communicate and document lessons."),
    ]
    for i, (t, b) in enumerate(steps):
        add_panel(s, 0.65 + i * 2.45, 2.0, 2.0, 1.25, t, b, accent=PPT_TEAL if i < 3 else PPT_RED)
        if i < 4:
            add_arrow(s, 2.62 + i * 2.45, 2.62, 3.05 + i * 2.45, 2.62, PPT_TEAL)
    add_textbox(s, 1.0, 4.5, 11.0, 0.7, "CERT-In directions require specified cyber incidents to be reported within 6 hours and ICT logs to be retained for 180 days within India.", size=19, bold=True, color=PPT_INK, align=PP_ALIGN.CENTER)

    s = new_slide("Certification", "The evidence pack is structured for formal review.")
    cert_points = [
        ("Privacy", "DPDP notice, consent, rights, retention and grievance artefacts."),
        ("Security", "VAPT, RLS tests, dependency scan, backup and incident drill."),
        ("Ethics", "IEC/SRC protocol, PIS/ICF and data management plan."),
        ("Operations", "Admin SOPs, training, support desk and monitoring."),
        ("Vendors", "Supabase, Vercel and AI-provider processor review."),
        ("Quality", "UAT, clinical safety review and release sign-off."),
    ]
    for i, (a, b) in enumerate(cert_points):
        x = 0.75 + (i % 3) * 4.0
        y = 1.55 + (i // 3) * 1.75
        add_panel(s, x, y, 3.45, 1.15, a, b, accent=PPT_BLUE if i % 2 else PPT_TEAL)

    s = new_slide("Patent Readiness", "The invention story is the FAP-specific AI and governance workflow.")
    add_panel(s, 0.8, 1.55, 3.2, 1.35, "Problem specificity", "Paper FAP logbooks cannot create auditable, longitudinal, AI-ready datasets.", accent=PPT_RED)
    add_panel(s, 4.3, 1.55, 3.2, 1.35, "Technical combination", "PWA field capture + RLS governance + AI reflection engine + mentor review + analytics.", accent=PPT_TEAL)
    add_panel(s, 7.8, 1.55, 3.2, 1.35, "Evidence trail", "Repo history, grant files, challenge submission, manuscripts, architecture and logs.", accent=PPT_GOLD)
    add_textbox(s, 1.1, 4.2, 10.7, 0.8, "Next IP step: prior-art search and provisional specification drafting around the integrated method, system architecture and AI-assisted educational assessment workflow.", size=20, bold=True, color=PPT_INK, align=PP_ALIGN.CENTER)

    s = new_slide("Scale Plan", "Full-scale implementation should proceed through controlled gates.")
    phases = [
        ("0", "Governance lock"),
        ("1", "Security validation"),
        ("2", "Pilot cohort"),
        ("3", "Institution scale"),
        ("4", "Research export"),
    ]
    for i, (num, label) in enumerate(phases):
        x = 0.9 + i * 2.4
        circle = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, PptInches(x), PptInches(2.0), PptInches(0.65), PptInches(0.65))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PPT_TEAL
        circle.line.fill.background()
        circle.text_frame.text = num
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        circle.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(255, 255, 255)
        circle.text_frame.paragraphs[0].font.bold = True
        add_textbox(s, x - 0.25, 2.85, 1.4, 0.55, label, size=13, bold=True, color=PPT_INK, align=PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            add_arrow(s, x + 0.75, 2.33, x + 2.25, 2.33, PPT_TEAL)
    add_bullet_box(s, 1.05, 4.25, 10.7, 0.9, [
        "Exit gate for full scale: signed governance pack, security sign-off, pilot report, trained admins/mentors and operational monitoring.",
    ], size=18)

    s = new_slide("Implementation Status", "The honest readiness message is strong, but should stay audit-defensible.")
    add_panel(s, 0.8, 1.55, 3.3, 1.3, "Ready to project", "Architecture, data model, RLS foundation, AI safety design, grant/research relevance and challenge recognition.", accent=PPT_TEAL)
    add_panel(s, 4.45, 1.55, 3.3, 1.3, "Ready to complete", "In-app consent, rights workflow, VAPT, RLS test automation, vendor agreements and incident drill.", accent=PPT_GOLD)
    add_panel(s, 8.1, 1.55, 3.3, 1.3, "Ready to scale", "After certification review and institutional sign-off, deploy cohort-wise with monitoring and support.", accent=PPT_BLUE)
    add_textbox(s, 1.2, 4.3, 10.3, 0.75, "Recommended wording: fully compliant-by-design and ready for patenting, certification review and full-scale implementation, subject to formal institutional/legal sign-offs.", size=21, bold=True, color=PPT_INK, align=PP_ALIGN.CENTER)

    s = new_slide("Sources", "The compliance pack is grounded in current official references.")
    add_bullet_box(
        s,
        0.8,
        1.45,
        11.8,
        4.9,
        [
            "MeitY: Digital Personal Data Protection Act, 2023.",
            "MeitY: Digital Personal Data Protection Rules, 2025 and enforcement timeline, notified November 2025.",
            "ICMR: National Ethical Guidelines for Biomedical and Health Research Involving Human Participants, 2017.",
            "CERT-In: Cyber Security Directions under IT Act section 70B, 28 April 2022.",
            "National Health Authority: ABDM overview and health-data governance direction.",
            "NMC: CBME / Family Adoption Programme public references.",
        ],
        size=17,
    )

    prs.core_properties.title = "FAP NextGen Architecture Compliance Certification Readiness"
    prs.core_properties.subject = "Architecture, security, DPDP compliance, certification and patent readiness"
    prs.core_properties.author = "FAP NextGen Project Team"
    prs.save(PPTX_PATH)


def main() -> None:
    build_figures()
    build_docx()
    build_pptx()
    print(DOCX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
